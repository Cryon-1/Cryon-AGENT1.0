#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import os
import argparse
from pptx import Presentation
import PyPDF2

def extract_ppt_text(ppt_path):
    try:
        prs = Presentation(ppt_path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
        return "\n".join(text_parts)
    except Exception as e:
        print(f"处理PPT {ppt_path} 时出错：{e}")
        return ""

def extract_pdf_text(pdf_path):
    """提取电子版PDF中的文本（非扫描）"""
    try:
        with open(pdf_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            text = []
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text.append(page_text)
            return "\n".join(text)
    except Exception as e:
        print(f"处理PDF {pdf_path} 时出错：{e}")
        return ""

def main():
    parser = argparse.ArgumentParser(description="提取PPT和PDF文件中的文本")
    parser.add_argument('--input_dir', required=True, help='存放文件的文件夹路径')
    parser.add_argument('--output_dir', required=True, help='输出文本文件的文件夹路径')
    args = parser.parse_args()

    input_dir = args.input_dir
    output_dir = args.output_dir

    if not os.path.isdir(input_dir):
        print(f"错误：输入文件夹不存在 - {input_dir}")
        return

    os.makedirs(output_dir, exist_ok=True)

    for filename in os.listdir(input_dir):
        file_path = os.path.join(input_dir, filename)
        if not os.path.isfile(file_path):
            continue

        ext = os.path.splitext(filename)[1].lower()
        base_name = os.path.splitext(filename)[0]
        out_txt = os.path.join(output_dir, base_name + '.txt')

        if os.path.exists(out_txt):
            print(f"跳过已存在文件：{out_txt}")
            continue

        text = ""
        if ext in ['.ppt', '.pptx']:
            print(f"正在处理PPT：{filename}")
            text = extract_ppt_text(file_path)
        elif ext == '.pdf':
            print(f"正在处理PDF：{filename}")
            text = extract_pdf_text(file_path)
        else:
            continue  # 跳过其他类型

        if text.strip():
            with open(out_txt, 'w', encoding='utf-8') as f:
                f.write(text)
            print(f"已保存：{out_txt}")
        else:
            print(f"警告：{filename} 未提取到文本")

    print("所有文件处理完成！")

if __name__ == "__main__":
    main()